# utils.py

import os
from groq import Groq
import PyPDF2
import docx
import pandas as pd
from pptx import Presentation
import re

def allowed_file(filename):
    return filename.lower().endswith((".mp4", ".mov", ".avi", ".mkv", ".wav", ".mp3", ".m4a", 
                                    ".pdf", ".doc", ".docx", ".txt", ".ppt", ".pptx", ".xls", ".xlsx"))


def chunk_text(text, chunk_size=500, overlap=50):
    """
    Split text into overlapping chunks for embedding.
    """
    chunks = []
    start = 0
    while start < len(text):
        end = min(start + chunk_size, len(text))
        chunks.append(text[start:end])
        start += chunk_size - overlap
    return chunks


def extract_text_from_file(file_path, file_type):
    """
    Extract text from different file types.
    """
    try:
        if file_type in ['pdf']:
            return extract_text_from_pdf(file_path)
        elif file_type in ['doc', 'docx']:
            return extract_text_from_docx(file_path)
        elif file_type in ['txt']:
            return extract_text_from_txt(file_path)
        elif file_type in ['ppt', 'pptx']:
            return extract_text_from_ppt(file_path)
        elif file_type in ['xls', 'xlsx']:
            return extract_text_from_excel(file_path)
        else:
            return f"[Error: Unsupported file type: {file_type}]"
    except Exception as e:
        return f"[Error extracting text from {file_type}: {e}]"


def extract_text_from_pdf(file_path):
    """Extract text from PDF file."""
    try:
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            text = ""
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
        return text.strip()
    except Exception as e:
        return f"[Error reading PDF: {e}]"


def extract_text_from_docx(file_path):
    """Extract text from DOCX file."""
    try:
        doc = docx.Document(file_path)
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        return text.strip()
    except Exception as e:
        return f"[Error reading DOCX: {e}]"


def extract_text_from_txt(file_path):
    """Extract text from TXT file."""
    try:
        with open(file_path, 'r', encoding='utf-8') as file:
            return file.read().strip()
    except Exception as e:
        return f"[Error reading TXT: {e}]"


def extract_text_from_ppt(file_path):
    """Extract text from PPT/PPTX file."""
    try:
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                # Try to extract text using getattr to avoid linter errors
                shape_text = getattr(shape, 'text', None)
                if shape_text:
                    text += shape_text + "\n"
                else:
                    # Try text_frame approach
                    text_frame = getattr(shape, 'text_frame', None)
                    if text_frame:
                        for paragraph in text_frame.paragraphs:
                            text += paragraph.text + "\n"
        return text.strip()
    except Exception as e:
        return f"[Error reading PPT: {e}]"


def extract_text_from_excel(file_path):
    """Extract text from Excel file."""
    try:
        # Read all sheets
        excel_file = pd.ExcelFile(file_path)
        text = ""
        for sheet_name in excel_file.sheet_names:
            df = pd.read_excel(file_path, sheet_name=sheet_name)
            text += f"Sheet: {sheet_name}\n"
            text += df.to_string() + "\n\n"
        return text.strip()
    except Exception as e:
        return f"[Error reading Excel: {e}]"


def get_file_type(filename):
    """Get file type from filename."""
    return filename.split('.')[-1].lower() if '.' in filename else None


def estimate_tokens_and_cost(text, model="llama3-70b-8192"):
    """
    Estimate token count and cost for Groq API calls.
    Rough estimation: 1 token ≈ 4 characters for English text.
    """
    # Rough token estimation (1 token ≈ 4 characters)
    estimated_tokens = len(text) // 4
    
    # Groq pricing (as of 2024) - approximate
    pricing = {
        "llama3-70b-8192": {"input": 0.00005, "output": 0.00010},  # per 1K tokens
        "llama3-8b-8192": {"input": 0.00002, "output": 0.00004},
        "mixtral-8x7b-32768": {"input": 0.00003, "output": 0.00006}
    }
    
    model_pricing = pricing.get(model, pricing["llama3-70b-8192"])
    
    # Estimate input and output tokens (assuming 1:1 ratio for metadata generation)
    input_tokens = estimated_tokens
    output_tokens = estimated_tokens * 0.8  # Output is usually shorter
    
    # Calculate cost
    input_cost = (input_tokens / 1000) * model_pricing["input"]
    output_cost = (output_tokens / 1000) * model_pricing["output"]
    total_cost = input_cost + output_cost
    
    return {
        "input_tokens": int(input_tokens),
        "output_tokens": int(output_tokens),
        "estimated_cost": total_cost
    }


def generate_video_metadata(transcript, groq_api_key=None):
    """
    Use Groq API to generate comprehensive video metadata from transcript.
    Returns a dictionary with title, description, highlights, takeaways, category, etc.
    """
    if groq_api_key is None:
        groq_api_key = os.getenv("GROK_API_KEY")
    if not groq_api_key:
        return {"error": "GROK_API_KEY not set"}
    
    # Estimate tokens and cost before making the call
    prompt = f"""
    Based on the following transcript, provide comprehensive metadata in JSON format:

    TRANSCRIPT:
    {transcript}

    Please analyze and provide the following information in a structured JSON format:
    {{
        "title": "A compelling, SEO-friendly title (max 60 characters)",
        "short_description": "A concise summary in 1-2 sentences (max 150 characters)",
        "detailed_description": "A comprehensive description of the content (2-3 paragraphs)",
        "key_highlights": ["Highlight 1", "Highlight 2", "Highlight 3", "Highlight 4", "Highlight 5"],
        "main_takeaways": ["Takeaway 1", "Takeaway 2", "Takeaway 3"],
        "category": "Primary category (e.g., Technology, Education, Entertainment, Business, etc.)",
        "subcategory": "More specific subcategory",
        "topics": ["Topic 1", "Topic 2", "Topic 3", "Topic 4"],
        "sentiment": "Overall sentiment (Positive, Negative, Neutral, Mixed)",
        "target_audience": "Who would benefit from this content",
        "estimated_duration": "Estimated duration in minutes based on content length",
        "difficulty_level": "Beginner, Intermediate, or Advanced",
        "action_items": ["Action item 1", "Action item 2", "Action item 3"],
        "related_concepts": ["Related concept 1", "Related concept 2", "Related concept 3"]
    }}

    Focus on accuracy and provide actionable insights. If the transcript is unclear or too short, indicate that in the response.
    """
    
    # Estimate cost
    cost_estimate = estimate_tokens_and_cost(prompt)
    
    try:
        client = Groq(api_key=groq_api_key)
        response = client.chat.completions.create(
            model="llama3-70b-8192",
            messages=[
                {"role": "system", "content": "You are an expert content analyst. Provide accurate, structured metadata in JSON format. Always respond with valid JSON."},
                {"role": "user", "content": prompt}
            ],
            max_tokens=1500,
            temperature=0.2
        )
        
        content = response.choices[0].message.content
        if content is None:
            return {"error": "No content returned from Groq API"}
        
        # Update cost estimate with actual usage
        actual_usage = {
            "input_tokens": getattr(response.usage, 'prompt_tokens', cost_estimate["input_tokens"]) if response.usage else cost_estimate["input_tokens"],
            "output_tokens": getattr(response.usage, 'completion_tokens', cost_estimate["output_tokens"]) if response.usage else cost_estimate["output_tokens"],
            "estimated_cost": cost_estimate["estimated_cost"]
        }
        
        # Try to extract JSON from the response
        import json
        try:
            # Find JSON in the response (sometimes it's wrapped in markdown)
            if "```json" in content:
                json_start = content.find("```json") + 7
                json_end = content.find("```", json_start)
                json_str = content[json_start:json_end].strip()
            elif "```" in content:
                json_start = content.find("```") + 3
                json_end = content.find("```", json_start)
                json_str = content[json_start:json_end].strip()
            else:
                json_str = content.strip()
            
            metadata = json.loads(json_str)
            metadata["token_usage"] = actual_usage  # Add token usage to metadata
            return metadata
            
        except json.JSONDecodeError:
            # If JSON parsing fails, return a structured response
            return {
                "title": "Content Analysis",
                "short_description": content[:150] + "..." if len(content) > 150 else content,
                "detailed_description": content,
                "key_highlights": ["Content analyzed successfully"],
                "main_takeaways": ["See detailed description"],
                "category": "General",
                "subcategory": "Analysis",
                "topics": ["Content Analysis"],
                "sentiment": "Neutral",
                "target_audience": "General",
                "estimated_duration": "Unknown",
                "difficulty_level": "General",
                "action_items": ["Review the detailed description"],
                "related_concepts": ["Content Analysis"],
                "token_usage": actual_usage
            }
            
    except Exception as e:
        return {"error": f"Groq API error: {e}"}
